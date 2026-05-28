"""PPTX text extraction using python-pptx."""

from __future__ import annotations

import io
from dataclasses import dataclass
from typing import Iterable

from pptx import Presentation

from shared.contracts import DEFAULT_CHUNK_TARGET_CHARS, EXTRACTION_STRATEGY_PPTX
from text_cleanup import normalize_text_block
from text_blocks import TextBlock, TextChunk, build_chunks_from_blocks


@dataclass(frozen=True)
class PptxExtractionResult:
    plain_text: str
    blocks: list[TextBlock]
    chunks: list[TextChunk]
    detected_languages: list[str]
    slide_count: int
    table_count: int
    extraction_strategy: str = EXTRACTION_STRATEGY_PPTX


def build_pptx_extraction(
    raw: bytes,
    *,
    chunk_target_chars: int = DEFAULT_CHUNK_TARGET_CHARS,
) -> PptxExtractionResult:
    presentation = Presentation(io.BytesIO(raw))
    blocks: list[TextBlock] = []
    table_count = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                table_count += 1
                table_text = _extract_table_text(shape.table)
                if table_text:
                    slide_parts.append(table_text)
                continue

            if getattr(shape, "has_text_frame", False):
                text = _extract_text_frame_text(shape.text_frame)
                if text:
                    slide_parts.append(text)

        slide_text = normalize_text_block("\n".join(slide_parts))
        if slide_text:
            blocks.append(
                TextBlock(
                    block_id=f"block-{len(blocks) + 1:06d}",
                    block_type="slide",
                    text=slide_text,
                    source_page_number=slide_number,
                )
            )

        notes_text = _extract_notes_text(slide)
        if notes_text:
            blocks.append(
                TextBlock(
                    block_id=f"block-{len(blocks) + 1:06d}",
                    block_type="speaker_notes",
                    text=notes_text,
                    source_page_number=slide_number,
                )
            )

    plain_text = "\n\n".join(block.text for block in blocks)
    chunks = build_chunks_from_blocks(blocks, chunk_target_chars=chunk_target_chars)

    return PptxExtractionResult(
        plain_text=plain_text,
        blocks=blocks,
        chunks=chunks,
        detected_languages=[],
        slide_count=len(presentation.slides),
        table_count=table_count,
    )


def _iter_shapes(shapes: Iterable) -> Iterable:
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _extract_text_frame_text(text_frame) -> str:
    paragraphs: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _extract_table_text(table) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _extract_notes_text(slide) -> str:
    # notes_slide creates a notes slide on access, so guard it first.
    if not slide.has_notes_slide:
        return ""
    text_frame = slide.notes_slide.notes_text_frame
    if text_frame is None:
        return ""
    return normalize_text_block(_extract_text_frame_text(text_frame))
