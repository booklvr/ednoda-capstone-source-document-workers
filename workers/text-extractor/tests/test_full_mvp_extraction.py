"""Unit tests for Full MVP CSV/DOCX/PPTX extraction and limit behavior."""

from __future__ import annotations

import sys
import unittest
import zipfile
from io import BytesIO
from pathlib import Path

WORKER_ROOT = Path(__file__).resolve().parents[2]
TEXT_EXTRACTOR_DIR = WORKER_ROOT / "text-extractor"
for path in (WORKER_ROOT, TEXT_EXTRACTOR_DIR):
    if str(path) not in sys.path:
        sys.path.insert(0, str(path))

from shared.contracts import (  # noqa: E402
    EXTRACTION_STATUS_PARTIAL,
    EXTRACTION_STATUS_READY,
    MANIFEST_VERSION,
    WARNING_EXTRACTION_BLOCKS_TRUNCATED,
    WARNING_EXTRACTION_CHARS_TRUNCATED,
)
from shared.keys import build_extraction_text_package_keys  # noqa: E402
from csv_logic import build_csv_extraction  # noqa: E402
from docx_logic import build_docx_extraction  # noqa: E402
from extraction_limits import apply_extraction_limits  # noqa: E402
from package_writer import build_manifest  # noqa: E402
from pptx_logic import build_pptx_extraction  # noqa: E402
from text_blocks import TextBlock  # noqa: E402
from txt_logic import build_txt_extraction  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402


def build_min_docx_bytes(*paragraphs: str) -> bytes:
    body_parts = []
    for paragraph in paragraphs:
        escaped = (
            paragraph.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )
        body_parts.append(
            f'<w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            f'<w:r><w:t>{escaped}</w:t></w:r></w:p>'
        )
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{''.join(body_parts)}</w:body></w:document>"
    )
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("word/document.xml", document_xml)
    return buffer.getvalue()


def build_pptx_bytes(
    *slides: str,
    tables: dict[int, list[list[str]]] | None = None,
    notes: dict[int, str] | None = None,
) -> bytes:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    tables = tables or {}
    notes = notes or {}

    for slide_index, slide_text in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        if slide_text:
            text_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.5),
                Inches(9),
                Inches(1.5),
            )
            text_box.text_frame.text = slide_text

        table_rows = tables.get(slide_index)
        if table_rows:
            row_count = len(table_rows)
            column_count = max(len(row) for row in table_rows)
            table_shape = slide.shapes.add_table(
                row_count,
                column_count,
                Inches(0.5),
                Inches(2),
                Inches(9),
                Inches(0.5 * row_count),
            )
            for row_index, row in enumerate(table_rows):
                for column_index, cell_text in enumerate(row):
                    table_shape.table.cell(row_index, column_index).text = cell_text

        note_text = notes.get(slide_index)
        if note_text:
            slide.notes_slide.notes_text_frame.text = note_text

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


class CsvExtractionTests(unittest.TestCase):
    def test_build_csv_extraction_preserves_headers_and_rows(self) -> None:
        raw = b"word,meaning\nhabitat,home\nmigrate,move seasonally\n"

        result = build_csv_extraction(raw)

        self.assertEqual(result.extraction_strategy, "csv")
        self.assertEqual(result.table_count, 1)
        self.assertEqual(len(result.blocks), 3)
        self.assertEqual(result.blocks[0].block_type, "table_header")
        self.assertEqual(result.blocks[1].block_type, "table_row")
        self.assertIn("word: habitat | meaning: home", result.plain_text)


class DocxExtractionTests(unittest.TestCase):
    def test_build_docx_extraction_preserves_paragraphs(self) -> None:
        raw = build_min_docx_bytes("Lesson paragraph one.", "Second paragraph.")

        result = build_docx_extraction(raw)

        self.assertEqual(result.extraction_strategy, "docx")
        self.assertGreaterEqual(len(result.blocks), 2)
        self.assertIn("Lesson paragraph one.", result.plain_text)
        self.assertIn("Second paragraph.", result.plain_text)

    def test_build_docx_extraction_preserves_repeated_lesson_lines(self) -> None:
        raw = build_min_docx_bytes(
            "Repeat after me.",
            "Photosynthesis: energy from light",
            "Repeat after me.",
            "What is chlorophyll?",
            "Repeat after me.",
            "Plants make sugar.",
        )

        result = build_docx_extraction(raw)

        self.assertEqual(result.plain_text.count("Repeat after me."), 3)
        self.assertIn("Photosynthesis: energy from light", result.plain_text)


class PptxExtractionTests(unittest.TestCase):
    def test_build_pptx_extraction_preserves_slide_boundaries(self) -> None:
        raw = build_pptx_bytes("Slide one title", "Slide two body")

        result = build_pptx_extraction(raw)

        self.assertEqual(result.extraction_strategy, "pptx")
        self.assertEqual(result.slide_count, 2)
        self.assertGreaterEqual(len(result.blocks), 2)
        self.assertIn("Slide one title", result.plain_text)
        self.assertIn("Slide two body", result.plain_text)

    def test_build_pptx_extraction_preserves_paragraph_boundaries(self) -> None:
        raw = build_pptx_bytes("Vocabulary\nphotosynthesis: energy from light")

        result = build_pptx_extraction(raw)

        self.assertIn("Vocabulary\nphotosynthesis", result.plain_text)
        self.assertNotIn("Vocabularyphotosynthesis", result.plain_text)

    def test_build_pptx_extraction_skips_empty_slides_but_keeps_slide_count(self) -> None:
        raw = build_pptx_bytes("", "Second slide")

        result = build_pptx_extraction(raw)

        self.assertEqual(result.slide_count, 2)
        self.assertEqual(len(result.blocks), 1)
        self.assertEqual(result.blocks[0].source_page_number, 2)
        self.assertGreaterEqual(len(result.chunks), 1)
        self.assertEqual(result.detected_languages, [])

    def test_build_pptx_extraction_extracts_tables_once_and_counts_them(self) -> None:
        raw = build_pptx_bytes(
            "Classroom vocabulary",
            tables={1: [["word", "meaning"], ["apple", "사과"], ["banana", "바나나"]]},
        )

        result = build_pptx_extraction(raw)

        self.assertEqual(result.table_count, 1)
        self.assertIn("Classroom vocabulary", result.plain_text)
        self.assertIn("word | meaning", result.plain_text)
        self.assertIn("apple | 사과", result.plain_text)
        self.assertEqual(result.plain_text.count("apple | 사과"), 1)

    def test_build_pptx_extraction_preserves_repeated_and_numbered_esl_content(self) -> None:
        raw = build_pptx_bytes(
            "1\napple\nRepeat after me.",
            "2\nbanana\nRepeat after me.",
            "3\norange\nRepeat after me.",
        )

        result = build_pptx_extraction(raw)

        self.assertEqual(result.plain_text.count("Repeat after me."), 3)
        self.assertIn("1\napple", result.plain_text)
        self.assertIn("2\nbanana", result.plain_text)
        self.assertIn("3\norange", result.plain_text)

    def test_build_pptx_extraction_attaches_existing_notes_to_correct_slide(self) -> None:
        raw = build_pptx_bytes(
            "Slide one",
            "Slide two",
            notes={2: "Teacher note for slide two"},
        )

        result = build_pptx_extraction(raw)

        notes_blocks = [block for block in result.blocks if block.block_type == "speaker_notes"]
        self.assertEqual(len(notes_blocks), 1)
        self.assertEqual(notes_blocks[0].source_page_number, 2)
        self.assertEqual(notes_blocks[0].text, "Teacher note for slide two")


class ExtractionLimitTests(unittest.TestCase):
    def test_chunking_splits_single_oversized_block(self) -> None:
        extracted = build_txt_extraction(
            b"First long sentence for chunk splitting. Second long sentence for chunk splitting.",
            chunk_target_chars=32,
        )

        self.assertGreater(len(extracted.chunks), 1)
        self.assertTrue(all(len(chunk.text) <= 40 for chunk in extracted.chunks))

    def test_apply_extraction_limits_marks_partial_when_blocks_truncated(self) -> None:
        extracted = build_txt_extraction(b"One.\n\nTwo.\n\nThree.")
        limited = apply_extraction_limits(
            plain_text=extracted.plain_text,
            blocks=extracted.blocks,
            chunks=extracted.chunks,
            max_chars=500_000,
            max_blocks=1,
            max_chunks=100,
        )

        self.assertEqual(limited.status, EXTRACTION_STATUS_PARTIAL)
        self.assertEqual(len(limited.blocks), 1)
        self.assertIn(WARNING_EXTRACTION_BLOCKS_TRUNCATED, limited.warnings)

    def test_apply_extraction_limits_marks_partial_when_chars_truncated(self) -> None:
        extracted = build_txt_extraction(b"Long lesson paragraph for truncation testing.")
        limited = apply_extraction_limits(
            plain_text=extracted.plain_text,
            blocks=extracted.blocks,
            chunks=extracted.chunks,
            max_chars=10,
            max_blocks=100,
            max_chunks=100,
        )

        self.assertEqual(limited.status, EXTRACTION_STATUS_PARTIAL)
        self.assertLessEqual(len(limited.plain_text), 10)
        self.assertIn(WARNING_EXTRACTION_CHARS_TRUNCATED, limited.warnings)


class ManifestContractTests(unittest.TestCase):
    def test_txt_manifest_validates_v1_contract_shape(self) -> None:
        raw = b"Alpha lesson notes.\n\nSecond paragraph."
        extracted = build_txt_extraction(raw)
        package_keys = build_extraction_text_package_keys(
            "ednoda-dev-source-document-text",
            "00000000-0000-4000-8000-000000000001",
            42,
            7,
        )
        manifest = build_manifest(
            source_document_id=42,
            extraction_id=7,
            original_filename="lesson.txt",
            original_mime_type="text/plain",
            file_extension=".txt",
            original_bucket="ednoda-dev-source-documents",
            original_key="source-documents/user/u/document/42/original/lesson.txt",
            text_bucket="ednoda-dev-source-document-text",
            package_keys=package_keys,
            extraction_strategy="plain_text",
            status=EXTRACTION_STATUS_READY,
            plain_text=extracted.plain_text,
            blocks=extracted.blocks,
            chunks=extracted.chunks,
            page_count=None,
            detected_languages=["en"],
            warnings=[],
        )

        self.assertEqual(manifest["version"], MANIFEST_VERSION)
        self.assertEqual(manifest["document"]["sourceDocumentId"], 42)
        self.assertEqual(manifest["outputs"]["plainText"]["key"], package_keys["plainTextKey"])


if __name__ == "__main__":
    unittest.main()
